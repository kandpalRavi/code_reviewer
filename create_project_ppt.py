from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


prs = Presentation()

# Theme colors
PRIMARY = RGBColor(24, 62, 125)
ACCENT = RGBColor(0, 150, 136)
DARK = RGBColor(44, 44, 44)
LIGHT_BG = RGBColor(245, 248, 252)


def add_title_box(slide, text, top=0.4):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.0), Inches(0.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = PRIMARY


def add_bullets(slide, items, start_top=1.5):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(start_top), Inches(11.5), Inches(5.3))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(12)
        p.font.size = Pt(22)
        p.font.color.rgb = DARK


def add_footer_bar(slide):
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(6.8), Inches(13.33), Inches(0.7)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = LIGHT_BG
    bar.line.fill.background()


# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
background.fill.solid()
background.fill.fore_color.rgb = LIGHT_BG
background.line.fill.background()

add_title_box(slide, "Code Reviewer", top=1.6)
subtitle = slide.shapes.add_textbox(Inches(0.9), Inches(2.6), Inches(11.8), Inches(1.4))
stf = subtitle.text_frame
p = stf.paragraphs[0]
p.text = "Automated Multi-Language Code Analysis Platform"
p.font.size = Pt(24)
p.font.color.rgb = ACCENT

# Bottom-left mentor and bottom-right student name
mentor = slide.shapes.add_textbox(Inches(0.6), Inches(6.6), Inches(4.8), Inches(0.5))
mtf = mentor.text_frame
mp = mtf.paragraphs[0]
mp.text = "Mentor: [Enter Mentor Name]"
mp.font.size = Pt(16)
mp.font.bold = True
mp.font.color.rgb = DARK

student = slide.shapes.add_textbox(Inches(8.3), Inches(6.6), Inches(4.4), Inches(0.5))
stf2 = student.text_frame
sp = stf2.paragraphs[0]
sp.text = "Presented by: Ravindra Kandpal"
sp.alignment = PP_ALIGN.RIGHT
sp.font.size = Pt(16)
sp.font.bold = True
sp.font.color.rgb = DARK

# Slide 2: Problem and Objective
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Problem Statement & Objective")
add_bullets(slide, [
    "Manual code reviews are slow and inconsistent.",
    "Developers need instant feedback on code quality and security.",
    "Objective: Build a full-stack tool to analyze Python, JavaScript, and Java code.",
    "Deliver actionable metrics, issues, and suggestions in one dashboard.",
])
add_footer_bar(slide)

# Slide 3: Architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "System Architecture")
add_bullets(slide, [
    "Frontend: React + Vite + Tailwind + Monaco Editor.",
    "Backend: FastAPI with modular analyzer services.",
    "Database: MongoDB for history, stats, and user data.",
    "Authentication: JWT tokens with secure password hashing.",
    "Pipeline: User input -> Analyzer -> Metrics + Issues -> UI Results.",
])
add_footer_bar(slide)

# Slide 4: Core Features
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Core Features Implemented")
add_bullets(slide, [
    "Real-time analysis for Python, JavaScript, and Java.",
    "Security scanning (Bandit), linting, and complexity checks.",
    "Maintainability score and issue severity classification.",
    "File upload, language auto-detection, and syntax highlighting.",
    "PDF export and analysis history tracking.",
])
add_footer_bar(slide)

# Slide 5: Technology Stack
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Technology Stack")
add_bullets(slide, [
    "Frontend: React 18, Vite, Tailwind CSS, Framer Motion, Recharts.",
    "Backend: Python, FastAPI, Pydantic, Motor (MongoDB).",
    "Analyzers: Pylint, Radon, Bandit, ESLint rules, Java static checks.",
    "Security: Passlib (bcrypt), Python-JOSE (JWT HS256).",
    "API: REST endpoints for analyze, auth, history, and stats.",
])
add_footer_bar(slide)

# Slide 6: Workflow Demo
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "User Workflow")
add_bullets(slide, [
    "1. User selects language or uploads source file.",
    "2. Code is submitted through /api/analyze endpoint.",
    "3. Backend runs analyzer chain and computes metrics.",
    "4. UI displays complexity, maintainability, and issue list.",
    "5. User can export report and review past analyses.",
])
add_footer_bar(slide)

# Slide 7: Testing and Results
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Testing & Validation Results")
add_bullets(slide, [
    "Unit Tests: 10/10 passed (Coverage: 88.5%, Time: 2.8s).",
    "System Tests: 10/10 passed (Avg Response: 450ms, Uptime: 100%).",
    "Verified key modules: auth, analyzers, database mapping, API responses.",
    "Validated edge cases: empty input, backend down, mobile responsiveness.",
])
add_footer_bar(slide)

# Slide 8: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Conclusion & Future Scope")
add_bullets(slide, [
    "Successfully delivered a robust code quality analysis platform.",
    "Improved developer productivity with fast and reliable feedback.",
    "Future scope: CI/CD integration, team dashboards, and custom rules.",
    "Future scope: LLM-assisted fix suggestions and multi-repo insights.",
    "Thank You",
])
add_footer_bar(slide)

output = "c:/Users/Ravindra Kandpal/Desktop/code_reviewer/Code_Reviewer_Project_Presentation.pptx"
prs.save(output)
print(f"Presentation created: {output}")
